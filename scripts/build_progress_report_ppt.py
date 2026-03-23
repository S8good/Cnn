import csv
import json
from pathlib import Path
from typing import Dict, List, Tuple

import matplotlib.pyplot as plt
import numpy as np
from matplotlib import patches
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_ROOT = PROJECT_ROOT / "outputs" / "report_assets_20260320_defense"
OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)


def choose_template() -> Path:
    candidates = sorted(PROJECT_ROOT.glob("*.pptx"))
    if not candidates:
        raise FileNotFoundError("No .pptx template found in project root.")
    for p in candidates:
        if "20260309" in p.name:
            return p
    return candidates[0]


PPT_TEMPLATE = choose_template()
PPT_OUTPUT = PROJECT_ROOT / "史鹏程 20260320_模型建设进展汇报_答辩风格.pptx"

# Thesis-defense style palette
NAVY = "#0B1D39"
BLUE = "#1E4E8C"
CYAN = "#1B7F8E"
ORANGE = "#C97824"
GREEN = "#2D9D72"
RED = "#C74B55"
TEXT_DARK = "#10243E"
TEXT_MUTED = "#5A6E86"
BG_SOFT = "#F5F8FC"
BORDER = "#D6DFEA"


def setup_matplotlib() -> None:
    plt.style.use("seaborn-v0_8-whitegrid")
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS"]
    plt.rcParams["axes.unicode_minus"] = False
    plt.rcParams["figure.facecolor"] = "white"
    plt.rcParams["axes.facecolor"] = "white"
    plt.rcParams["savefig.bbox"] = "tight"
    plt.rcParams["savefig.pad_inches"] = 0.15
    plt.rcParams["axes.titleweight"] = "bold"
    plt.rcParams["axes.labelsize"] = 11
    plt.rcParams["axes.titlesize"] = 13
    plt.rcParams["xtick.labelsize"] = 10
    plt.rcParams["ytick.labelsize"] = 10


def read_json(path: Path) -> Dict:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def read_csv_rows(path: Path) -> List[Dict[str, str]]:
    with path.open("r", encoding="utf-8", newline="") as f:
        return list(csv.DictReader(f))


def load_train_metrics() -> Dict:
    return read_json(PROJECT_ROOT / "outputs" / "run_20260320_085725" / "train_metrics.json")


def load_fewshot_summary() -> List[Dict[str, str]]:
    return read_csv_rows(PROJECT_ROOT / "outputs" / "fewshot_bundle_cea_v1" / "bundle_summary.csv")


def load_real_dataset() -> Tuple[np.ndarray, np.ndarray, np.ndarray, Dict[int, str]]:
    data_dir = PROJECT_ROOT / "data" / "real_fewshot_cea"
    spectra = np.load(data_dir / "spectra.npy").astype(np.float32)
    labels = np.load(data_dir / "labels.npy").astype(np.int64)
    wavelengths = np.load(data_dir / "wavelengths.npy").astype(np.float32)
    label_map = read_json(data_dir / "label_map.json")["id_to_label"]
    id_to_label = {int(k): str(v) for k, v in label_map.items()}
    return spectra, labels, wavelengths, id_to_label


def build_pretrain_asset(path: Path, train_metrics: Dict) -> None:
    history = train_metrics["history"]
    epochs = np.asarray([int(x["epoch"]) for x in history], dtype=np.int64)
    train_loss = np.asarray([float(x["train_loss"]) for x in history], dtype=np.float64)
    val_loss = np.asarray([float(x["val_loss"]) for x in history], dtype=np.float64)
    val_acc = np.asarray([float(x["val_acc"]) for x in history], dtype=np.float64) * 100.0

    fig, axes = plt.subplots(1, 2, figsize=(12.8, 4.6), dpi=240)

    axes[0].plot(epochs, train_loss, marker="o", lw=2.4, color=BLUE, label="train loss")
    axes[0].plot(epochs, val_loss, marker="s", lw=2.4, color=ORANGE, label="val loss")
    axes[0].fill_between(epochs, 0, train_loss, color=BLUE, alpha=0.08)
    axes[0].set_xlabel("Epoch")
    axes[0].set_ylabel("Loss")
    axes[0].set_title("预训练损失曲线")
    axes[0].legend(frameon=False)
    axes[0].grid(alpha=0.3, linestyle="--")

    axes[1].plot(epochs, val_acc, marker="o", lw=2.6, color=CYAN)
    best_idx = int(np.argmax(val_acc))
    axes[1].scatter([epochs[best_idx]], [val_acc[best_idx]], color=RED, s=60, zorder=5)
    axes[1].annotate(
        f"Best {val_acc[best_idx]:.2f}%",
        xy=(epochs[best_idx], val_acc[best_idx]),
        xytext=(epochs[best_idx] - 2.2, val_acc[best_idx] - 13),
        arrowprops=dict(arrowstyle="->", color=RED, lw=1.4),
        fontsize=10.5,
        color=RED,
    )
    axes[1].set_xlabel("Epoch")
    axes[1].set_ylabel("Validation Accuracy (%)")
    axes[1].set_title("预训练验证精度")
    axes[1].set_ylim(0, 102)
    axes[1].grid(alpha=0.3, linestyle="--")

    best = train_metrics["best"]
    fig.suptitle(
        f"Stage A 预训练结果：val_acc={best['val_acc']*100:.2f}%  val_loss={best['val_loss']:.3f}",
        fontsize=15.5,
        color=TEXT_DARK,
        fontweight="bold",
    )
    fig.savefig(path)
    plt.close(fig)


def build_real_data_asset(path: Path, spectra: np.ndarray, labels: np.ndarray, wavelengths: np.ndarray, id_to_label: Dict[int, str]) -> None:
    uniq, cnt = np.unique(labels, return_counts=True)
    order = np.argsort(cnt)[::-1]
    classes = [id_to_label[int(uniq[i])] for i in order]
    counts = [int(cnt[i]) for i in order]

    fig, axes = plt.subplots(1, 2, figsize=(12.8, 4.8), dpi=240)

    colors = [BLUE if c >= 6 else ORANGE for c in counts]
    bars = axes[0].bar(np.arange(len(classes)), counts, color=colors, alpha=0.92, edgecolor="white", linewidth=0.6)
    axes[0].axhline(4, color=CYAN, linestyle="--", lw=1.8, label="k=3,n_query=1 下限")
    axes[0].axhline(6, color=ORANGE, linestyle="--", lw=1.8, label="k=5,n_query=1 下限")
    axes[0].set_xticks(np.arange(len(classes)))
    axes[0].set_xticklabels(classes, rotation=40, ha="right")
    axes[0].set_ylabel("样本数")
    axes[0].set_title("真实数据类别分布")
    axes[0].legend(frameon=False, fontsize=9.5, loc="upper left")
    axes[0].grid(axis="y", alpha=0.3, linestyle="--")
    for bar, val in zip(bars, counts):
        axes[0].text(bar.get_x() + bar.get_width() / 2, val + 0.35, str(val), ha="center", va="bottom", fontsize=9.2, color=TEXT_DARK)

    focus_ids = [0, 2, 4, 7, 8, 10]
    cmap = plt.colormaps["viridis"].resampled(len(focus_ids))
    for idx, cid in enumerate(focus_ids):
        sample_idx = np.where(labels == cid)[0]
        if sample_idx.size == 0:
            continue
        mean_curve = spectra[sample_idx].mean(axis=0)
        std_curve = spectra[sample_idx].std(axis=0)
        color = cmap(idx)
        axes[1].plot(wavelengths, mean_curve, lw=2.0, color=color, label=id_to_label[cid])
        axes[1].fill_between(wavelengths, mean_curve - std_curve, mean_curve + std_curve, color=color, alpha=0.14)

    axes[1].set_title("不同浓度的平均光谱与波动范围")
    axes[1].set_xlabel("Wavelength (nm)")
    axes[1].set_ylabel("Normalized delta signal")
    axes[1].legend(frameon=False, fontsize=8.8, ncol=2)
    axes[1].grid(alpha=0.3, linestyle="--")

    fig.suptitle("Stage B 真实数据概览：168 样本 / 11 类 / 400 采样点（类别不均衡）", fontsize=15.5, color=TEXT_DARK, fontweight="bold")
    fig.savefig(path)
    plt.close(fig)


def build_fewshot_asset(path: Path, rows: List[Dict[str, str]]) -> None:
    valid = [r for r in rows if r.get("status") == "ok"]
    grouped = {"prototype": [], "linear_head": []}
    for r in valid:
        grouped[r["mode"]].append(r)
    for mode in grouped:
        grouped[mode].sort(key=lambda x: int(x["k_shot"]))

    fig, axes = plt.subplots(1, 2, figsize=(12.8, 4.8), dpi=240)
    mode_color = {"prototype": BLUE, "linear_head": ORANGE}

    for mode, vals in grouped.items():
        if not vals:
            continue
        ks = np.asarray([int(v["k_shot"]) for v in vals], dtype=np.int64)
        acc = np.asarray([float(v["acc_mean"]) * 100 for v in vals], dtype=np.float64)
        acc_std = np.asarray([float(v["acc_std"]) * 100 for v in vals], dtype=np.float64)
        f1 = np.asarray([float(v["macro_f1_mean"]) * 100 for v in vals], dtype=np.float64)
        f1_std = np.asarray([float(v["macro_f1_std"]) * 100 for v in vals], dtype=np.float64)

        axes[0].errorbar(ks, acc, yerr=acc_std, marker="o", capsize=4, lw=2.3, color=mode_color[mode], label=mode)
        axes[1].errorbar(ks, f1, yerr=f1_std, marker="o", capsize=4, lw=2.3, color=mode_color[mode], label=mode)

    random_baseline = {1: 100 / 11, 3: 100 / 9, 5: 100 / 8}
    axes[0].plot(list(random_baseline.keys()), list(random_baseline.values()), linestyle="--", color=TEXT_MUTED, lw=1.8, label="random baseline")

    axes[0].set_title("Few-shot Query Accuracy")
    axes[1].set_title("Few-shot Macro-F1")
    for ax in axes:
        ax.set_xlabel("K-shot")
        ax.set_ylabel("Score (%)")
        ax.set_xticks([1, 3, 5])
        ax.set_ylim(0, 40)
        ax.grid(alpha=0.3, linestyle="--")
        ax.legend(frameon=False, fontsize=9.5)

    fig.suptitle("Stage B few-shot 结果：已高于随机基线，但整体表现仍偏低", fontsize=15.5, color=TEXT_DARK, fontweight="bold")
    fig.savefig(path)
    plt.close(fig)


def build_two_step_architecture_asset(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(13.6, 4.9), dpi=240)
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 40)
    ax.axis("off")

    ax.add_patch(patches.Rectangle((0, 0), 100, 40, facecolor="#F8FAFD", edgecolor="none"))

    # Bottom foundation: Step 1
    ax.add_patch(
        patches.FancyBboxPatch(
            (8, 7),
            84,
            11,
            boxstyle="round,pad=0.8,rounding_size=2.2",
            linewidth=2.0,
            edgecolor=BLUE,
            facecolor="#EDF4FD",
        )
    )
    ax.text(12, 15.3, "Step 1  Meta-Baseline 特征预训练（底座）", fontsize=16.5, color=BLUE, fontweight="bold", va="center")
    ax.text(
        12,
        11.2,
        "理论光谱网格化 -> 物理噪声增强 -> Whole-classification 预训练 -> 保留 128 维 Encoder",
        fontsize=11.6,
        color=TEXT_DARK,
        va="center",
    )
    ax.add_patch(
        patches.FancyBboxPatch(
            (81.5, 13.0),
            8.2,
            3.0,
            boxstyle="round,pad=0.35,rounding_size=0.8",
            linewidth=0,
            facecolor=GREEN,
        )
    )
    ax.text(85.6, 14.5, "已完成", fontsize=10.8, color="white", fontweight="bold", ha="center", va="center")

    # Step 2 amplifier
    ax.add_patch(
        patches.FancyBboxPatch(
            (26, 22),
            48,
            8,
            boxstyle="round,pad=0.8,rounding_size=2.2",
            linewidth=2.0,
            edgecolor=ORANGE,
            facecolor="#FFF6EC",
        )
    )
    ax.text(30, 27.1, "Step 2  cGAN 稀缺数据扩增（放大器）", fontsize=15.2, color=ORANGE, fontweight="bold", va="center")
    ax.text(30, 24.0, "聚焦低浓度 / LOD 附近 / 过渡区样本缺口，做受控条件生成与物理筛选", fontsize=11.2, color=TEXT_DARK, va="center")
    ax.add_patch(
        patches.FancyBboxPatch(
            (66.0, 24.4),
            6.0,
            2.3,
            boxstyle="round,pad=0.25,rounding_size=0.7",
            linewidth=0,
            facecolor=ORANGE,
        )
    )
    ax.text(69.0, 25.55, "规划中", fontsize=9.8, color="white", fontweight="bold", ha="center", va="center")

    # Top outcome
    ax.add_patch(
        patches.FancyBboxPatch(
            (34, 33),
            32,
            3.5,
            boxstyle="round,pad=0.4,rounding_size=1.4",
            linewidth=1.6,
            edgecolor=NAVY,
            facecolor="white",
        )
    )
    ax.text(50, 34.75, "目标：更稳的真实 few-shot 适配与更强的低浓度鲁棒性", fontsize=11.4, color=NAVY, fontweight="bold", ha="center", va="center")

    ax.annotate("", xy=(50, 22), xytext=(50, 18), arrowprops=dict(arrowstyle="->", lw=2.2, color=TEXT_MUTED))
    ax.annotate("", xy=(50, 33), xytext=(50, 30), arrowprops=dict(arrowstyle="->", lw=2.2, color=TEXT_MUTED))

    ax.text(4, 36.5, "两步关系：先建立物理锚点，再按需放大低浓度脆弱区", fontsize=16.5, color=TEXT_DARK, fontweight="bold")
    ax.text(4, 3.0, "没有 Step 1 的物理锚点，Step 2 只会生成“看起来像光谱”的伪数据。", fontsize=12.0, color=RED)

    fig.savefig(path)
    plt.close(fig)


def build_step1_pipeline_asset(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(14.2, 4.9), dpi=240)
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 26)
    ax.axis("off")

    ax.add_patch(patches.Rectangle((0, 0), 100, 26, facecolor="#F8FAFD", edgecolor="none"))

    items = [
        ("物理参数\n离散化", "粒径 / 折射率\n连续空间切网格", BLUE),
        ("批量生成\n理论库", "Mie / FDTD\n完美无噪声光谱", CYAN),
        ("物理噪声\n注入增强", "高斯噪声 / 基线漂移\n峰宽展宽", GREEN),
        ("Whole-classification\n预训练", "用分类任务压榨\n表征学习能力", ORANGE),
        ("弃用分类头\n保留 Encoder", "留下真正有价值的\n128 维表征空间", NAVY),
    ]

    x0 = 2.5
    for idx, (title, body, color) in enumerate(items):
        x = x0 + idx * 19.0
        ax.add_patch(
            patches.FancyBboxPatch(
                (x, 7.0),
                15.8,
                11.4,
                boxstyle="round,pad=0.55,rounding_size=1.6",
                linewidth=1.8,
                edgecolor=color,
                facecolor="white",
            )
        )
        ax.text(x + 1.0, 16.9, title, fontsize=12.8, color=color, fontweight="bold", va="top")
        ax.text(x + 1.0, 8.8, body, fontsize=10.6, color=TEXT_DARK, va="bottom", linespacing=1.45)
        if idx < len(items) - 1:
            ax.annotate("", xy=(x + 18.4, 12.7), xytext=(x + 15.9, 12.7), arrowprops=dict(arrowstyle="->", lw=2.1, color=TEXT_MUTED))

    ax.text(2.5, 22.2, "阶段一目标：逼迫模型学出稳定、抗干扰的光谱高维表征空间，而不是急于做最终检测。", fontsize=15.8, color=TEXT_DARK, fontweight="bold")
    ax.text(2.5, 2.5, "本质：先学 physical manifold，再做真实域迁移。", fontsize=11.8, color=TEXT_MUTED)

    fig.savefig(path)
    plt.close(fig)


def build_step2_pipeline_asset(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(14.0, 4.9), dpi=240)
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 30)
    ax.axis("off")

    ax.add_patch(patches.Rectangle((0, 0), 100, 30, facecolor="#F8FAFD", edgecolor="none"))

    items = [
        (4, 9, 18, 11, BLUE, "1. 锁定稀缺区", "只针对低浓度 / LOD / 过渡区\n拒绝全范围盲目扩增"),
        (27, 9, 20, 11, CYAN, "2. 构建条件输入", "基础谱形 + 目标条件\n而不是纯噪声乱生成"),
        (52, 9, 18, 11, ORANGE, "3. 学习受控映射", "从已知状态到微弱状态\n保持物理连续性"),
        (75, 9, 20, 11, GREEN, "4. 物理规则筛选", "峰位 / 峰宽 / 单调性通过后\n才允许进入训练库"),
    ]

    for i, (x, y, w, h, color, title, body) in enumerate(items):
        ax.add_patch(
            patches.FancyBboxPatch(
                (x, y),
                w,
                h,
                boxstyle="round,pad=0.6,rounding_size=1.8",
                linewidth=1.8,
                edgecolor=color,
                facecolor="white",
            )
        )
        ax.text(x + 1.1, y + h - 2.0, title, fontsize=13.3, color=color, fontweight="bold", va="top")
        ax.text(x + 1.1, y + 2.0, body, fontsize=10.8, color=TEXT_DARK, va="bottom", linespacing=1.45)
        if i < len(items) - 1:
            ax.annotate("", xy=(x + w + 3.5, 14.5), xytext=(x + w + 0.4, 14.5), arrowprops=dict(arrowstyle="->", lw=2.1, color=TEXT_MUTED))

    ax.add_patch(
        patches.FancyBboxPatch(
            (33.5, 22.7),
            33.0,
            3.1,
            boxstyle="round,pad=0.35,rounding_size=0.9",
            linewidth=0,
            facecolor="#FCEDEE",
        )
    )
    ax.text(50, 24.25, "边界提醒：cGAN 是“受控增强”，不是“凭空造谱”。", fontsize=11.6, color=RED, fontweight="bold", ha="center", va="center")

    ax.text(4, 26.9, "阶段二目标：补齐最贵、最稀缺的低浓度样本缺口，稳住 few-shot 原型锚点。", fontsize=15.4, color=TEXT_DARK, fontweight="bold")

    fig.savefig(path)
    plt.close(fig)


def generate_assets() -> Dict[str, Path]:
    setup_matplotlib()
    train_metrics = load_train_metrics()
    fewshot_rows = load_fewshot_summary()
    spectra, labels, wavelengths, id_to_label = load_real_dataset()

    assets = {
        "architecture": OUTPUT_ROOT / "two_step_architecture.png",
        "step1_pipeline": OUTPUT_ROOT / "step1_pipeline.png",
        "step2_pipeline": OUTPUT_ROOT / "step2_pipeline.png",
        "pretrain": OUTPUT_ROOT / "pretrain_summary.png",
        "real_data": OUTPUT_ROOT / "real_data_overview.png",
        "fewshot": OUTPUT_ROOT / "fewshot_summary.png",
        "tsne": PROJECT_ROOT / "outputs" / "run_20260320_085725" / "tsne_validation.png",
    }

    build_two_step_architecture_asset(assets["architecture"])
    build_step1_pipeline_asset(assets["step1_pipeline"])
    build_step2_pipeline_asset(assets["step2_pipeline"])
    build_pretrain_asset(assets["pretrain"], train_metrics)
    build_real_data_asset(assets["real_data"], spectra, labels, wavelengths, id_to_label)
    build_fewshot_asset(assets["fewshot"], fewshot_rows)
    return assets


def rgb(hex_code: str) -> RGBColor:
    raw = hex_code.replace("#", "")
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def remove_all_slides(prs: Presentation) -> None:
    for idx in range(len(prs.slides) - 1, -1, -1):
        slide_id = prs.slides._sldIdLst[idx]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[idx]


def set_title(slide, text: str) -> None:
    shape = slide.shapes.title
    shape.text = text
    p = shape.text_frame.paragraphs[0]
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = rgb(TEXT_DARK)


def add_text_card(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=14,
    color=TEXT_DARK,
    bold=False,
    align=PP_ALIGN.LEFT,
    fill=None,
    line=BORDER,
    radius=True,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1.0)
    else:
        shape.line.color.rgb = rgb("FFFFFF")
        shape.line.transparency = 1.0

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for line_text in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = line_text
        p.alignment = align
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.bold = bold if first else False
        p.font.color.rgb = rgb(color)
        first = False
    return shape


def add_header_tag(slide, text: str) -> None:
    add_text_card(slide, Cm(24.3), Cm(0.45), Cm(6.0), Cm(0.95), text, font_size=9.5, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, fill=BLUE, line=BLUE)


def add_footer(slide, page_no: int) -> None:
    add_text_card(slide, Cm(0.8), Cm(18.1), Cm(22.0), Cm(0.8), "LSPR 模型建设进展汇报", font_size=9.5, color=TEXT_MUTED, line=None, radius=False)
    add_text_card(slide, Cm(31.0), Cm(18.0), Cm(1.0), Cm(0.9), str(page_no), font_size=10, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, fill=BLUE, line=BLUE)


def add_picture(slide, path: Path, left, top, width, height) -> None:
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.background()
    frame.line.color.rgb = rgb(BORDER)
    frame.line.width = Pt(1.1)


def add_bullets_card(slide, left, top, width, height, title: str, bullets: List[str], fill=BG_SOFT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(BORDER)
    shape.line.width = Pt(1.0)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = "Microsoft YaHei"
    p0.font.size = Pt(15.5)
    p0.font.bold = True
    p0.font.color.rgb = rgb(TEXT_DARK)

    for line in bullets:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.bullet = True
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12.2)
        p.font.color.rgb = rgb(TEXT_MUTED)


def set_cell_text(cell, text: str, size: float = 11.0, bold: bool = False, color: str = TEXT_DARK, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    cell.text = ""
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)


def add_progress_table(slide, left, top, width, height) -> None:
    rows = [
        ("A1", "理论光谱生成", "已实现", "Mie 光谱生成与光学常数读取已完成"),
        ("A2", "物理噪声增强", "已实现", "高斯噪声、基线漂移、峰位/FWHM 扰动已完成"),
        ("A3", "Encoder 预训练", "已实现", "1D-ResNet 预训练完成，最佳 val_acc=98.91%"),
        ("A4", "表征验证", "已实现", "t-SNE 与关键指标已输出"),
        ("B1", "真实数据导入", "已实现", "CSV 与成对 Excel 数据已可直接处理"),
        ("B2", "few-shot 适配", "已实现", "prototype 与 linear_head 均已跑通"),
        ("B3", "多 episode 评估", "已实现", "K=1/3/5 汇总实验已完成"),
        ("B4", "单样本推理", "已实现", "离线单样本推理流程已完成"),
        ("C1", "物理约束判定", "待实现", "红移方向与边界约束尚未接入"),
        ("C2", "可信度输出", "待实现", "校准与不确定性解释尚未形成"),
        ("D1", "软件集成", "待实现", "尚未接入目标测试软件入口"),
        ("D2", "自动报告闭环", "待实现", "日志数据库与报告流程未完成"),
    ]

    headers = ["编码", "模块", "状态", "说明"]
    table_shape = slide.shapes.add_table(len(rows) + 1, 4, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Cm(2.0)
    table.columns[1].width = Cm(5.0)
    table.columns[2].width = Cm(3.0)
    table.columns[3].width = Cm(19.8)

    header_fill = rgb(NAVY)
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        set_cell_text(cell, head, size=12, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

    for r, row in enumerate(rows, start=1):
        bg = "#FBFDFF" if r % 2 == 1 else "#FFFFFF"
        for c in range(4):
            table.cell(r, c).fill.solid()
            table.cell(r, c).fill.fore_color.rgb = rgb(bg)

        set_cell_text(table.cell(r, 0), row[0], size=11.2, bold=True, align=PP_ALIGN.CENTER)
        set_cell_text(table.cell(r, 1), row[1], size=11.2)

        status = row[2]
        status_cell = table.cell(r, 2)
        if status == "已实现":
            status_cell.fill.solid()
            status_cell.fill.fore_color.rgb = rgb(GREEN)
            set_cell_text(status_cell, status, size=10.8, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)
        else:
            status_cell.fill.solid()
            status_cell.fill.fore_color.rgb = rgb(ORANGE)
            set_cell_text(status_cell, status, size=10.8, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

        set_cell_text(table.cell(r, 3), row[3], size=11.0, color=TEXT_MUTED)


def add_content_slide(prs: Presentation, title: str, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(slide, title)
    add_header_tag(slide, "LSPR 答辩汇报")
    add_footer(slide, page_no)
    return slide


def add_title_slide(prs: Presentation, assets: Dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_text_card(slide, Cm(0.9), Cm(0.7), Cm(12.0), Cm(0.9), "LSPR 答辩汇报", font_size=9.8, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, fill=BLUE, line=BLUE)
    add_text_card(
        slide,
        Cm(1.0),
        Cm(2.8),
        Cm(17.8),
        Cm(3.5),
        "融合物理先验与条件生成的少样本 LSPR 光谱分析架构",
        font_size=25.5,
        color=TEXT_DARK,
        bold=True,
        line=None,
        radius=False,
    )
    add_text_card(
        slide,
        Cm(1.0),
        Cm(6.6),
        Cm(16.0),
        Cm(1.9),
        "从理论预训练到极限数据扩增的深度学习策略",
        font_size=15.5,
        color=TEXT_MUTED,
        bold=False,
        line=None,
        radius=False,
    )
    add_text_card(
        slide,
        Cm(1.0),
        Cm(10.0),
        Cm(8.8),
        Cm(1.4),
        "汇报人：史鹏程",
        font_size=13.8,
        color=TEXT_MUTED,
        bold=True,
        line=None,
        radius=False,
    )
    add_text_card(
        slide,
        Cm(1.0),
        Cm(11.2),
        Cm(8.8),
        Cm(1.2),
        "日期：2026.03.22",
        font_size=12.8,
        color=TEXT_MUTED,
        line=None,
        radius=False,
    )
    add_text_card(
        slide,
        Cm(1.0),
        Cm(14.4),
        Cm(17.0),
        Cm(2.2),
        "答辩视角：先讲清“为什么必须先做物理底座”，\n再说明“为什么低浓度区需要条件生成放大器”。",
        font_size=13.5,
        color=TEXT_DARK,
        fill="#ECF3FB",
        line=BORDER,
    )
    add_picture(slide, assets["tsne"], Cm(19.6), Cm(2.6), Cm(11.1), Cm(11.1))
    add_text_card(
        slide,
        Cm(19.8),
        Cm(14.3),
        Cm(10.7),
        Cm(2.5),
        "视觉锚点：预训练后的特征流形已经显现出结构化分簇，\n说明 128 维 Encoder 具备“物理感知”能力。",
        font_size=12.0,
        color=TEXT_DARK,
        fill="#FFF7EE",
        line=BORDER,
    )
    add_footer(slide, 1)


def build_presentation(assets: Dict[str, Path]) -> None:
    prs = Presentation(str(PPT_TEMPLATE))
    remove_all_slides(prs)

    add_title_slide(prs, assets)

    slide = add_content_slide(prs, "LSPR 深度学习应用的核心痛点", 2)
    add_bullets_card(
        slide,
        Cm(0.8),
        Cm(3.1),
        Cm(11.0),
        Cm(6.0),
        "核心矛盾",
        [
            "真实 LSPR 光谱获取成本高、周期长，深度模型长期处于“数据饥饿”状态。",
            "高浓度样本相对容易获得，低浓度 / 临界响应区样本最有价值却最稀缺。",
            "直接端到端训练容易死记少量曲线，缺乏对峰位、峰宽、基线扰动规律的物理理解。",
        ],
    )
    add_text_card(
        slide,
        Cm(0.8),
        Cm(9.6),
        Cm(11.0),
        Cm(4.2),
        "方案引出\n\n我们需要一个两层架构：\n先建立物理常识底座，再对最稀缺的低浓度区做受控放大。",
        font_size=14.0,
        color=TEXT_DARK,
        fill="#F4FAF9",
        line=BORDER,
    )
    add_picture(slide, assets["real_data"], Cm(12.2), Cm(3.0), Cm(18.7), Cm(11.1))
    add_text_card(
        slide,
        Cm(12.5),
        Cm(14.2),
        Cm(12.0),
        Cm(1.1),
        "实测数据证据：168 样本 / 11 类 / 400 采样点，且类别明显不均衡。",
        font_size=11.8,
        color="#FFFFFF",
        bold=True,
        align=PP_ALIGN.CENTER,
        fill=BLUE,
        line=BLUE,
    )

    slide = add_content_slide(prs, "两步递进式解决方案：底座与放大器", 3)
    add_picture(slide, assets["architecture"], Cm(0.9), Cm(2.9), Cm(30.3), Cm(9.8))
    add_text_card(
        slide,
        Cm(1.0),
        Cm(13.1),
        Cm(14.5),
        Cm(2.4),
        "Step 1（必须做）\nMeta-Baseline 特征预训练：建立 128 维物理表征空间。",
        font_size=13.5,
        color=TEXT_DARK,
        fill="#ECF3FB",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(15.8),
        Cm(13.1),
        Cm(14.8),
        Cm(2.4),
        "Step 2（进阶做）\ncGAN 稀缺数据扩增：补足 LOD 附近最贵、最稀缺的样本缺口。",
        font_size=13.5,
        color=TEXT_DARK,
        fill="#FFF7EE",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(1.0),
        Cm(15.9),
        Cm(29.6),
        Cm(1.5),
        "当前项目状态：Step 1 与真实 few-shot 验证已完成；Step 2 尚未实现，但已被当前结果明确地“推出来”。",
        font_size=12.6,
        color=TEXT_DARK,
        fill="#F7FAFD",
        line=BORDER,
    )

    slide = add_content_slide(prs, "阶段一：基于理论光谱的物理感知预训练", 4)
    add_picture(slide, assets["step1_pipeline"], Cm(0.9), Cm(3.0), Cm(30.3), Cm(8.6))
    add_text_card(
        slide,
        Cm(1.0),
        Cm(12.0),
        Cm(10.5),
        Cm(4.2),
        "核心目标\n\n不急于做最终检测，而是先逼 encoder 学出稳定、抗干扰的高维表征空间。",
        font_size=13.8,
        color=TEXT_DARK,
        fill="#ECF3FB",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(11.9),
        Cm(12.0),
        Cm(9.0),
        Cm(4.2),
        "为什么是 Whole-classification\n\n表面上做分类，实质上是在压榨表征学习能力，为真实 few-shot 迁移做地基。",
        font_size=13.2,
        color=TEXT_DARK,
        fill="#F4FAF9",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(21.4),
        Cm(12.0),
        Cm(9.6),
        Cm(4.2),
        "当前对应成果\n\nMie 理论光谱生成、物理噪声增强、1D-ResNet 预训练与 t-SNE 验证均已完成。",
        font_size=13.0,
        color=TEXT_DARK,
        fill="#FFF7EE",
        line=BORDER,
    )

    slide = add_content_slide(prs, "阶段一：从“死记曲线”到“理解物理流形”", 5)
    add_picture(slide, assets["pretrain"], Cm(0.9), Cm(3.0), Cm(11.7), Cm(5.8))
    add_picture(slide, assets["tsne"], Cm(13.0), Cm(3.0), Cm(17.8), Cm(8.6))
    add_text_card(
        slide,
        Cm(1.0),
        Cm(9.2),
        Cm(11.5),
        Cm(3.5),
        "关键结果\n\n最佳验证精度：98.91%\n最佳验证损失：0.062\n表征维度：128",
        font_size=13.4,
        color=TEXT_DARK,
        fill="#F4FAF9",
        line=BORDER,
    )
    add_bullets_card(
        slide,
        Cm(1.0),
        Cm(13.2),
        Cm(14.5),
        Cm(4.0),
        "学术意义",
        [
            "先学习 physical manifold，再做真实域迁移。",
            "真实数据只需承担域适配，而不必从头教模型认识光谱。",
            "清晰分簇说明嵌入空间已具有结构化物理表征能力。",
        ],
    )
    add_text_card(
        slide,
        Cm(16.0),
        Cm(13.2),
        Cm(14.7),
        Cm(4.0),
        "边界提醒\n\n这证明的是“仿真域表征成功”，不是“真实域任务已经完成”。它是地基，不是终点。",
        font_size=13.0,
        color=TEXT_DARK,
        fill="#FFF7EE",
        line=BORDER,
    )

    slide = add_content_slide(prs, "阶段二：利用 cGAN 攻克极限数据缺口", 6)
    add_picture(slide, assets["step2_pipeline"], Cm(0.9), Cm(3.0), Cm(30.3), Cm(8.8))
    add_text_card(
        slide,
        Cm(1.0),
        Cm(12.2),
        Cm(14.8),
        Cm(3.8),
        "四步实操重点\n\n靶向低浓度缺口、构建条件输入、学习受控映射、通过物理规则筛选后再入库。",
        font_size=13.2,
        color=TEXT_DARK,
        fill="#ECF3FB",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(16.2),
        Cm(12.2),
        Cm(14.6),
        Cm(3.8),
        "边界强调\n\ncGAN 不是替代实验，也不是替代 Step 1；它是建立在物理锚点基础上的 V2.0 放大器。",
        font_size=13.2,
        color=TEXT_DARK,
        fill="#FFF7EE",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(23.8),
        Cm(1.1),
        Cm(6.4),
        Cm(0.9),
        "当前状态：规划中",
        font_size=9.8,
        color="FFFFFF",
        bold=True,
        align=PP_ALIGN.CENTER,
        fill=ORANGE,
        line=ORANGE,
    )

    slide = add_content_slide(prs, "为什么第二步值得做：来自当前结果的直接证据", 7)
    add_picture(slide, assets["fewshot"], Cm(0.9), Cm(3.0), Cm(18.7), Cm(10.9))
    add_bullets_card(
        slide,
        Cm(19.9),
        Cm(3.0),
        Cm(11.0),
        Cm(5.7),
        "当前观察",
        [
            "真实 few-shot 已高于随机基线，但整体 accuracy / macro-F1 仍偏低。",
            "support 集偏小、类别不均衡、低浓度样本稀缺，会直接导致 prototype 锚点漂移。",
            "这说明当前瓶颈不在“流程没跑通”，而在“关键区间样本不够稳”。",
        ],
        fill="#F7FAFD",
    )
    add_bullets_card(
        slide,
        Cm(19.9),
        Cm(9.1),
        Cm(11.0),
        Cm(5.2),
        "Step 2 的潜在收益",
        [
            "补齐 LOD 附近最贵、最稀缺的数据空白。",
            "增强低浓度与过渡区的鲁棒性。",
            "帮助探索极限检测附近的非线性物理响应。",
        ],
        fill="#FFF7EE",
    )
    add_text_card(
        slide,
        Cm(20.0),
        Cm(14.8),
        Cm(10.9),
        Cm(1.5),
        "当前结论：Step 1 已落地，Step 2 不是“可有可无”，而是下一阶段最自然的增强方向。",
        font_size=12.2,
        color=TEXT_DARK,
        fill="#ECF3FB",
        line=BORDER,
    )

    slide = add_content_slide(prs, "总结：物理先验与条件生成的完美闭环", 8)
    add_text_card(
        slide,
        Cm(1.0),
        Cm(3.0),
        Cm(29.8),
        Cm(1.8),
        "核心结论：两步架构不是并列替代，而是严格的递进关系。",
        font_size=17.2,
        color=TEXT_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
        fill="#ECF3FB",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(1.0),
        Cm(5.5),
        Cm(9.4),
        Cm(5.9),
        "Take-home 1\n\n底座先行：Step 1 先把“数字游标卡尺”做出来，让模型先理解 LSPR 光谱的物理结构。",
        font_size=13.2,
        color=TEXT_DARK,
        fill="#F7FAFD",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(11.2),
        Cm(5.5),
        Cm(9.4),
        Cm(5.9),
        "Take-home 2\n\n按需放大：Step 2 只在低浓度 / 弱响应区明显短缺时接入，目标是补齐最贵的短板。",
        font_size=13.2,
        color=TEXT_DARK,
        fill="#FFF7EE",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(21.4),
        Cm(5.5),
        Cm(9.4),
        Cm(5.9),
        "优先级建议\n\nMVP：巩固 Meta-Baseline + 真实 few-shot。\nV2.0：围绕低浓度样本缺口接入 cGAN 放大器。",
        font_size=13.0,
        color=TEXT_DARK,
        fill="#F4FAF9",
        line=BORDER,
    )
    add_text_card(
        slide,
        Cm(1.0),
        Cm(12.3),
        Cm(29.8),
        Cm(3.3),
        "Take-home message：没有物理锚点的 GAN 只会制造伪数据；有了物理底座，条件生成才能成为真正可用的低浓度放大器。",
        font_size=15.0,
        color=TEXT_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
        fill="#ECF3FB",
        line=BORDER,
    )

    prs.save(str(PPT_OUTPUT))


def main() -> None:
    assets = generate_assets()
    build_presentation(assets)
    print(f"Template: {PPT_TEMPLATE}")
    print(f"Assets: {OUTPUT_ROOT}")
    print(f"Saved PPT: {PPT_OUTPUT}")


if __name__ == "__main__":
    main()
